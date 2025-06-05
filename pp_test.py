from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import os

# Couleurs personnalisées
BLUE_LIGHT = RGBColor(173, 216, 230)  # bleu clair
BLACK = RGBColor(0, 0, 0)
WHITE = RGBColor(255, 255, 255)

# Création de la présentation
prs = Presentation()

# Chemin vers le logo
LOGO_PATH = "logo_nuage.png"

# Fonction pour ajouter une diapositive stylée
def add_custom_slide(title_text, body_text):
    slide_layout = prs.slide_layouts[1]  # Titre + contenu
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    content = slide.placeholders[1]

    # Style titre
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.bold = True
    title.text_frame.paragraphs[0].font.color.rgb = BLACK

    # Style texte
    content.text = body_text
    for paragraph in content.text_frame.paragraphs:
        paragraph.font.size = Pt(20)
        paragraph.font.color.rgb = BLACK

    # Ajouter logo
    if os.path.exists(LOGO_PATH):
        slide.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.0))

    # Optionnel : fond bleu clair (limité avec python-pptx)
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = WHITE  # Fond blanc pour rester propre

# Slide 1 – Titre
slide_title = prs.slides.add_slide(prs.slide_layouts[0])
slide_title.shapes.title.text = "Utilisation de Google Dialogflow dans une PME"
slide_title.placeholders[1].text = "Présentation et démo – Juin 2025"
# Ajouter logo
if os.path.exists(LOGO_PATH):
    slide_title.shapes.add_picture(LOGO_PATH, Inches(9), Inches(0.2), width=Inches(1.0))

# Slides de contenu
add_custom_slide("Contexte", "PME souhaitant améliorer l’interaction client.\nExplorer le potentiel des chatbots cloud avec Google Dialogflow.")
add_custom_slide("Qu’est-ce que Dialogflow ?", "- Plateforme Google Cloud pour créer des chatbots\n- Basée sur le NLP\n- Version ES (standard) et CX (avancée)")
add_custom_slide("Cas d’usage PME", "- Service client\n- Réservations\n- FAQ automatisée\n- Assistance RH ou IT")
add_custom_slide("Démo : Chatbot FAQ", "- 3 intents : horaires, disponibilité, contact\n- Créé dans Dialogflow ES\n- Réponses personnalisées simples")
add_custom_slide("Architecture simplifiée", "- Agent Dialogflow → Webhook facultatif\n- Intégration via iframe dans site PME\n- Backend possible via Firebase")
add_custom_slide("Coût estimé", "- 1 000 requêtes/jour gratuites\n- 0.002 USD / requête au-delà\n- Exemple : 60 000 requêtes = ~60 USD / mois")
add_custom_slide("Avantages & Limites", "✅ Rapide à déployer, NLP Google\n❌ Vendor lock-in, peu de contrôle\n❌ Pas open source comme Rasa")
add_custom_slide("Quand utiliser / éviter", "👍 Utiliser pour support client, FAQ, projets simples\n👎 Éviter si besoin de contrôle local ou dialogue complexe")
add_custom_slide("Ressources utiles", "- https://cloud.google.com/dialogflow/docs\n- Codelabs, tutoriels, formations")

# Enregistre le fichier
prs.save("Dialogflow_PME_Presentation.pptx")
print("✅ Présentation créée : Dialogflow_PME_Presentation.pptx")
